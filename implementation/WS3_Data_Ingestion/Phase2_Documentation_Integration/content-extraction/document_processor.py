"""
Document Processor for Nexus Architect
Comprehensive multi-format document processing with intelligent content extraction
"""

import os
import re
import json
import logging
import time
import asyncio
import hashlib
import tempfile
import mimetypes
from typing import Dict, List, Optional, Any, Set, Tuple, Union, BinaryIO
from dataclasses import dataclass, field
from enum import Enum
from datetime import datetime
from pathlib import Path
import aiofiles
import aiohttp
from PIL import Image
import pytesseract
from prometheus_client import Counter, Histogram, Gauge

# Document processing libraries
try:
    import docx
    from docx.document import Document as DocxDocument
except ImportError:
    docx = None

try:
    import pptx
    from pptx import Presentation
except ImportError:
    pptx = None

try:
    import openpyxl
    from openpyxl import load_workbook
except ImportError:
    openpyxl = None

try:
    import PyPDF2
    import fitz  # PyMuPDF
except ImportError:
    PyPDF2 = None
    fitz = None

try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None

try:
    import markdown
    from markdown.extensions import toc, tables, codehilite
except ImportError:
    markdown = None

# NLP libraries
try:
    import spacy
    from spacy import displacy
except ImportError:
    spacy = None

try:
    import nltk
    from nltk.tokenize import sent_tokenize, word_tokenize
    from nltk.corpus import stopwords
    from nltk.stem import WordNetLemmatizer
except ImportError:
    nltk = None

# Metrics
DOCUMENT_PROCESSING_REQUESTS = Counter('document_processing_requests_total', 'Total document processing requests', ['format', 'status'])
DOCUMENT_PROCESSING_LATENCY = Histogram('document_processing_latency_seconds', 'Document processing latency', ['format'])
CONTENT_EXTRACTION_ACCURACY = Gauge('content_extraction_accuracy', 'Content extraction accuracy score', ['format'])
OCR_PROCESSING_TIME = Histogram('ocr_processing_time_seconds', 'OCR processing time', ['image_type'])

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class ProcessingStatus(Enum):
    """Document processing status"""
    PENDING = "pending"
    PROCESSING = "processing"
    COMPLETED = "completed"
    FAILED = "failed"
    PARTIAL = "partial"

class ContentType(Enum):
    """Types of extracted content"""
    TEXT = "text"
    HEADING = "heading"
    PARAGRAPH = "paragraph"
    LIST = "list"
    TABLE = "table"
    IMAGE = "image"
    CODE = "code"
    LINK = "link"
    METADATA = "metadata"

@dataclass
class ExtractedContent:
    """Extracted content with metadata"""
    content_type: ContentType
    text: str
    metadata: Dict[str, Any] = field(default_factory=dict)
    position: Optional[int] = None
    confidence: float = 1.0
    language: Optional[str] = None

@dataclass
class DocumentStructure:
    """Document structure information"""
    headings: List[Dict[str, Any]] = field(default_factory=list)
    sections: List[Dict[str, Any]] = field(default_factory=list)
    tables: List[Dict[str, Any]] = field(default_factory=list)
    images: List[Dict[str, Any]] = field(default_factory=list)
    links: List[Dict[str, Any]] = field(default_factory=list)
    footnotes: List[Dict[str, Any]] = field(default_factory=list)
    toc: List[Dict[str, Any]] = field(default_factory=list)

@dataclass
class ProcessingResult:
    """Document processing result"""
    document_id: str
    format: str
    status: ProcessingStatus
    extracted_text: str
    structured_content: List[ExtractedContent]
    document_structure: DocumentStructure
    metadata: Dict[str, Any]
    processing_time: float
    error_message: Optional[str] = None
    confidence_score: float = 1.0
    language: Optional[str] = None
    word_count: int = 0
    character_count: int = 0

class PDFProcessor:
    """PDF document processor with OCR support"""
    
    def __init__(self):
        self.ocr_enabled = self._check_ocr_availability()
    
    def _check_ocr_availability(self) -> bool:
        """Check if OCR is available"""
        try:
            pytesseract.get_tesseract_version()
            return True
        except Exception:
            logger.warning("Tesseract OCR not available")
            return False
    
    async def process_pdf(self, file_path: str) -> ProcessingResult:
        """Process PDF document"""
        start_time = time.time()
        document_id = hashlib.md5(file_path.encode()).hexdigest()
        
        try:
            extracted_text = ""
            structured_content = []
            document_structure = DocumentStructure()
            metadata = {}
            
            # Try PyMuPDF first (better for complex PDFs)
            if fitz:
                result = await self._process_with_pymupdf(file_path)
                if result:
                    extracted_text, structured_content, document_structure, metadata = result
            
            # Fallback to PyPDF2
            if not extracted_text and PyPDF2:
                result = await self._process_with_pypdf2(file_path)
                if result:
                    extracted_text, structured_content, document_structure, metadata = result
            
            # OCR fallback for image-based PDFs
            if not extracted_text.strip() and self.ocr_enabled:
                result = await self._process_with_ocr(file_path)
                if result:
                    extracted_text, structured_content, document_structure, metadata = result
            
            processing_time = time.time() - start_time
            
            # Calculate confidence based on extraction success
            confidence = 1.0 if extracted_text.strip() else 0.5
            
            result = ProcessingResult(
                document_id=document_id,
                format="pdf",
                status=ProcessingStatus.COMPLETED if extracted_text else ProcessingStatus.PARTIAL,
                extracted_text=extracted_text,
                structured_content=structured_content,
                document_structure=document_structure,
                metadata=metadata,
                processing_time=processing_time,
                confidence_score=confidence,
                word_count=len(extracted_text.split()),
                character_count=len(extracted_text)
            )
            
            DOCUMENT_PROCESSING_REQUESTS.labels(format='pdf', status='success').inc()
            CONTENT_EXTRACTION_ACCURACY.labels(format='pdf').set(confidence)
            
            return result
        
        except Exception as e:
            processing_time = time.time() - start_time
            logger.error(f"PDF processing failed: {e}")
            
            DOCUMENT_PROCESSING_REQUESTS.labels(format='pdf', status='error').inc()
            
            return ProcessingResult(
                document_id=document_id,
                format="pdf",
                status=ProcessingStatus.FAILED,
                extracted_text="",
                structured_content=[],
                document_structure=DocumentStructure(),
                metadata={},
                processing_time=processing_time,
                error_message=str(e)
            )
        
        finally:
            latency = time.time() - start_time
            DOCUMENT_PROCESSING_LATENCY.labels(format='pdf').observe(latency)
    
    async def _process_with_pymupdf(self, file_path: str) -> Optional[Tuple]:
        """Process PDF with PyMuPDF"""
        try:
            doc = fitz.open(file_path)
            extracted_text = ""
            structured_content = []
            document_structure = DocumentStructure()
            metadata = doc.metadata
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                
                # Extract text
                page_text = page.get_text()
                extracted_text += page_text + "\n"
                
                # Extract structured content
                blocks = page.get_text("dict")
                for block in blocks.get("blocks", []):
                    if "lines" in block:
                        for line in block["lines"]:
                            for span in line.get("spans", []):
                                text = span.get("text", "").strip()
                                if text:
                                    # Determine content type based on font size and style
                                    font_size = span.get("size", 12)
                                    font_flags = span.get("flags", 0)
                                    
                                    if font_size > 16 or font_flags & 2**4:  # Large or bold
                                        content_type = ContentType.HEADING
                                    else:
                                        content_type = ContentType.TEXT
                                    
                                    structured_content.append(ExtractedContent(
                                        content_type=content_type,
                                        text=text,
                                        metadata={
                                            "page": page_num + 1,
                                            "font_size": font_size,
                                            "font_flags": font_flags,
                                            "bbox": span.get("bbox")
                                        }
                                    ))
                
                # Extract images
                image_list = page.get_images()
                for img_index, img in enumerate(image_list):
                    xref = img[0]
                    pix = fitz.Pixmap(doc, xref)
                    
                    document_structure.images.append({
                        "page": page_num + 1,
                        "index": img_index,
                        "width": pix.width,
                        "height": pix.height,
                        "colorspace": pix.colorspace.name if pix.colorspace else None
                    })
                    
                    pix = None  # Free memory
                
                # Extract links
                links = page.get_links()
                for link in links:
                    document_structure.links.append({
                        "page": page_num + 1,
                        "uri": link.get("uri"),
                        "bbox": link.get("from"),
                        "type": "external" if link.get("uri", "").startswith("http") else "internal"
                    })
            
            doc.close()
            return extracted_text, structured_content, document_structure, metadata
        
        except Exception as e:
            logger.error(f"PyMuPDF processing failed: {e}")
            return None
    
    async def _process_with_pypdf2(self, file_path: str) -> Optional[Tuple]:
        """Process PDF with PyPDF2"""
        try:
            extracted_text = ""
            structured_content = []
            document_structure = DocumentStructure()
            metadata = {}
            
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                # Extract metadata
                if pdf_reader.metadata:
                    metadata = {
                        "title": pdf_reader.metadata.get("/Title"),
                        "author": pdf_reader.metadata.get("/Author"),
                        "subject": pdf_reader.metadata.get("/Subject"),
                        "creator": pdf_reader.metadata.get("/Creator"),
                        "producer": pdf_reader.metadata.get("/Producer"),
                        "creation_date": pdf_reader.metadata.get("/CreationDate"),
                        "modification_date": pdf_reader.metadata.get("/ModDate")
                    }
                
                # Extract text from each page
                for page_num, page in enumerate(pdf_reader.pages):
                    try:
                        page_text = page.extract_text()
                        extracted_text += page_text + "\n"
                        
                        # Create structured content
                        if page_text.strip():
                            structured_content.append(ExtractedContent(
                                content_type=ContentType.TEXT,
                                text=page_text,
                                metadata={"page": page_num + 1}
                            ))
                    
                    except Exception as e:
                        logger.warning(f"Failed to extract text from page {page_num + 1}: {e}")
            
            return extracted_text, structured_content, document_structure, metadata
        
        except Exception as e:
            logger.error(f"PyPDF2 processing failed: {e}")
            return None
    
    async def _process_with_ocr(self, file_path: str) -> Optional[Tuple]:
        """Process PDF with OCR"""
        if not self.ocr_enabled:
            return None
        
        try:
            start_time = time.time()
            
            doc = fitz.open(file_path)
            extracted_text = ""
            structured_content = []
            document_structure = DocumentStructure()
            metadata = {"ocr_processed": True}
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                
                # Convert page to image
                mat = fitz.Matrix(2, 2)  # 2x zoom for better OCR
                pix = page.get_pixmap(matrix=mat)
                img_data = pix.tobytes("png")
                
                # Save to temporary file for OCR
                with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_file:
                    temp_file.write(img_data)
                    temp_path = temp_file.name
                
                try:
                    # Perform OCR
                    ocr_text = pytesseract.image_to_string(temp_path, lang='eng')
                    extracted_text += ocr_text + "\n"
                    
                    if ocr_text.strip():
                        structured_content.append(ExtractedContent(
                            content_type=ContentType.TEXT,
                            text=ocr_text,
                            metadata={
                                "page": page_num + 1,
                                "ocr_confidence": 0.8  # Estimated
                            },
                            confidence=0.8
                        ))
                
                finally:
                    # Clean up temporary file
                    os.unlink(temp_path)
                
                pix = None  # Free memory
            
            doc.close()
            
            ocr_time = time.time() - start_time
            OCR_PROCESSING_TIME.labels(image_type='pdf').observe(ocr_time)
            
            return extracted_text, structured_content, document_structure, metadata
        
        except Exception as e:
            logger.error(f"OCR processing failed: {e}")
            return None

class OfficeDocumentProcessor:
    """Microsoft Office document processor"""
    
    async def process_docx(self, file_path: str) -> ProcessingResult:
        """Process DOCX document"""
        start_time = time.time()
        document_id = hashlib.md5(file_path.encode()).hexdigest()
        
        try:
            if not docx:
                raise ImportError("python-docx library not available")
            
            doc = docx.Document(file_path)
            extracted_text = ""
            structured_content = []
            document_structure = DocumentStructure()
            
            # Extract metadata
            metadata = {
                "title": doc.core_properties.title,
                "author": doc.core_properties.author,
                "subject": doc.core_properties.subject,
                "keywords": doc.core_properties.keywords,
                "created": doc.core_properties.created,
                "modified": doc.core_properties.modified,
                "last_modified_by": doc.core_properties.last_modified_by
            }
            
            # Process paragraphs
            for para_index, paragraph in enumerate(doc.paragraphs):
                text = paragraph.text.strip()
                if not text:
                    continue
                
                extracted_text += text + "\n"
                
                # Determine content type based on style
                style_name = paragraph.style.name.lower()
                if "heading" in style_name:
                    content_type = ContentType.HEADING
                    level = self._extract_heading_level(style_name)
                    document_structure.headings.append({
                        "level": level,
                        "text": text,
                        "position": para_index
                    })
                else:
                    content_type = ContentType.PARAGRAPH
                
                structured_content.append(ExtractedContent(
                    content_type=content_type,
                    text=text,
                    metadata={
                        "style": paragraph.style.name,
                        "position": para_index
                    },
                    position=para_index
                ))
            
            # Process tables
            for table_index, table in enumerate(doc.tables):
                table_data = []
                for row in table.rows:
                    row_data = []
                    for cell in row.cells:
                        row_data.append(cell.text.strip())
                    table_data.append(row_data)
                
                document_structure.tables.append({
                    "index": table_index,
                    "rows": len(table.rows),
                    "columns": len(table.columns) if table.rows else 0,
                    "data": table_data
                })
                
                # Add table as structured content
                table_text = "\n".join(["\t".join(row) for row in table_data])
                extracted_text += table_text + "\n"
                
                structured_content.append(ExtractedContent(
                    content_type=ContentType.TABLE,
                    text=table_text,
                    metadata={
                        "table_index": table_index,
                        "rows": len(table.rows),
                        "columns": len(table.columns) if table.rows else 0
                    }
                ))
            
            processing_time = time.time() - start_time
            
            result = ProcessingResult(
                document_id=document_id,
                format="docx",
                status=ProcessingStatus.COMPLETED,
                extracted_text=extracted_text,
                structured_content=structured_content,
                document_structure=document_structure,
                metadata=metadata,
                processing_time=processing_time,
                confidence_score=1.0,
                word_count=len(extracted_text.split()),
                character_count=len(extracted_text)
            )
            
            DOCUMENT_PROCESSING_REQUESTS.labels(format='docx', status='success').inc()
            CONTENT_EXTRACTION_ACCURACY.labels(format='docx').set(1.0)
            
            return result
        
        except Exception as e:
            processing_time = time.time() - start_time
            logger.error(f"DOCX processing failed: {e}")
            
            DOCUMENT_PROCESSING_REQUESTS.labels(format='docx', status='error').inc()
            
            return ProcessingResult(
                document_id=document_id,
                format="docx",
                status=ProcessingStatus.FAILED,
                extracted_text="",
                structured_content=[],
                document_structure=DocumentStructure(),
                metadata={},
                processing_time=processing_time,
                error_message=str(e)
            )
        
        finally:
            latency = time.time() - start_time
            DOCUMENT_PROCESSING_LATENCY.labels(format='docx').observe(latency)
    
    async def process_pptx(self, file_path: str) -> ProcessingResult:
        """Process PPTX presentation"""
        start_time = time.time()
        document_id = hashlib.md5(file_path.encode()).hexdigest()
        
        try:
            if not pptx:
                raise ImportError("python-pptx library not available")
            
            presentation = Presentation(file_path)
            extracted_text = ""
            structured_content = []
            document_structure = DocumentStructure()
            
            # Extract metadata
            metadata = {
                "title": presentation.core_properties.title,
                "author": presentation.core_properties.author,
                "subject": presentation.core_properties.subject,
                "created": presentation.core_properties.created,
                "modified": presentation.core_properties.modified,
                "slide_count": len(presentation.slides)
            }
            
            # Process slides
            for slide_index, slide in enumerate(presentation.slides):
                slide_text = f"Slide {slide_index + 1}:\n"
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        slide_text += text + "\n"
                        
                        # Determine if it's a title or content
                        if hasattr(shape, "placeholder_format") and shape.placeholder_format:
                            if shape.placeholder_format.type == 1:  # Title placeholder
                                content_type = ContentType.HEADING
                                document_structure.headings.append({
                                    "level": 1,
                                    "text": text,
                                    "slide": slide_index + 1
                                })
                            else:
                                content_type = ContentType.TEXT
                        else:
                            content_type = ContentType.TEXT
                        
                        structured_content.append(ExtractedContent(
                            content_type=content_type,
                            text=text,
                            metadata={
                                "slide": slide_index + 1,
                                "shape_type": str(shape.shape_type)
                            }
                        ))
                
                extracted_text += slide_text + "\n"
            
            processing_time = time.time() - start_time
            
            result = ProcessingResult(
                document_id=document_id,
                format="pptx",
                status=ProcessingStatus.COMPLETED,
                extracted_text=extracted_text,
                structured_content=structured_content,
                document_structure=document_structure,
                metadata=metadata,
                processing_time=processing_time,
                confidence_score=1.0,
                word_count=len(extracted_text.split()),
                character_count=len(extracted_text)
            )
            
            DOCUMENT_PROCESSING_REQUESTS.labels(format='pptx', status='success').inc()
            CONTENT_EXTRACTION_ACCURACY.labels(format='pptx').set(1.0)
            
            return result
        
        except Exception as e:
            processing_time = time.time() - start_time
            logger.error(f"PPTX processing failed: {e}")
            
            DOCUMENT_PROCESSING_REQUESTS.labels(format='pptx', status='error').inc()
            
            return ProcessingResult(
                document_id=document_id,
                format="pptx",
                status=ProcessingStatus.FAILED,
                extracted_text="",
                structured_content=[],
                document_structure=DocumentStructure(),
                metadata={},
                processing_time=processing_time,
                error_message=str(e)
            )
        
        finally:
            latency = time.time() - start_time
            DOCUMENT_PROCESSING_LATENCY.labels(format='pptx').observe(latency)
    
    async def process_xlsx(self, file_path: str) -> ProcessingResult:
        """Process XLSX spreadsheet"""
        start_time = time.time()
        document_id = hashlib.md5(file_path.encode()).hexdigest()
        
        try:
            if not openpyxl:
                raise ImportError("openpyxl library not available")
            
            workbook = load_workbook(file_path, data_only=True)
            extracted_text = ""
            structured_content = []
            document_structure = DocumentStructure()
            
            # Extract metadata
            metadata = {
                "title": workbook.properties.title,
                "author": workbook.properties.creator,
                "subject": workbook.properties.subject,
                "created": workbook.properties.created,
                "modified": workbook.properties.modified,
                "sheet_count": len(workbook.worksheets),
                "sheet_names": workbook.sheetnames
            }
            
            # Process worksheets
            for sheet_index, worksheet in enumerate(workbook.worksheets):
                sheet_text = f"Sheet: {worksheet.title}\n"
                
                # Extract data from cells
                sheet_data = []
                for row in worksheet.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        row_data = [str(cell) if cell is not None else "" for cell in row]
                        sheet_data.append(row_data)
                        sheet_text += "\t".join(row_data) + "\n"
                
                document_structure.tables.append({
                    "sheet_name": worksheet.title,
                    "sheet_index": sheet_index,
                    "rows": len(sheet_data),
                    "columns": len(sheet_data[0]) if sheet_data else 0,
                    "data": sheet_data
                })
                
                structured_content.append(ExtractedContent(
                    content_type=ContentType.TABLE,
                    text=sheet_text,
                    metadata={
                        "sheet_name": worksheet.title,
                        "sheet_index": sheet_index,
                        "rows": len(sheet_data),
                        "columns": len(sheet_data[0]) if sheet_data else 0
                    }
                ))
                
                extracted_text += sheet_text + "\n"
            
            processing_time = time.time() - start_time
            
            result = ProcessingResult(
                document_id=document_id,
                format="xlsx",
                status=ProcessingStatus.COMPLETED,
                extracted_text=extracted_text,
                structured_content=structured_content,
                document_structure=document_structure,
                metadata=metadata,
                processing_time=processing_time,
                confidence_score=1.0,
                word_count=len(extracted_text.split()),
                character_count=len(extracted_text)
            )
            
            DOCUMENT_PROCESSING_REQUESTS.labels(format='xlsx', status='success').inc()
            CONTENT_EXTRACTION_ACCURACY.labels(format='xlsx').set(1.0)
            
            return result
        
        except Exception as e:
            processing_time = time.time() - start_time
            logger.error(f"XLSX processing failed: {e}")
            
            DOCUMENT_PROCESSING_REQUESTS.labels(format='xlsx', status='error').inc()
            
            return ProcessingResult(
                document_id=document_id,
                format="xlsx",
                status=ProcessingStatus.FAILED,
                extracted_text="",
                structured_content=[],
                document_structure=DocumentStructure(),
                metadata={},
                processing_time=processing_time,
                error_message=str(e)
            )
        
        finally:
            latency = time.time() - start_time
            DOCUMENT_PROCESSING_LATENCY.labels(format='xlsx').observe(latency)
    
    def _extract_heading_level(self, style_name: str) -> int:
        """Extract heading level from style name"""
        match = re.search(r'heading\s*(\d+)', style_name, re.IGNORECASE)
        return int(match.group(1)) if match else 1

class MarkdownProcessor:
    """Markdown document processor"""
    
    async def process_markdown(self, file_path: str) -> ProcessingResult:
        """Process Markdown document"""
        start_time = time.time()
        document_id = hashlib.md5(file_path.encode()).hexdigest()
        
        try:
            async with aiofiles.open(file_path, 'r', encoding='utf-8') as file:
                content = await file.read()
            
            extracted_text = content
            structured_content = []
            document_structure = DocumentStructure()
            
            # Extract metadata from front matter
            metadata = self._extract_front_matter(content)
            
            # Remove front matter from content
            content = self._remove_front_matter(content)
            
            # Parse markdown structure
            lines = content.split('\n')
            current_position = 0
            
            for line_num, line in enumerate(lines):
                line = line.strip()
                if not line:
                    continue
                
                # Extract headings
                if line.startswith('#'):
                    level = len(line) - len(line.lstrip('#'))
                    text = line.lstrip('#').strip()
                    
                    document_structure.headings.append({
                        "level": level,
                        "text": text,
                        "line": line_num + 1
                    })
                    
                    structured_content.append(ExtractedContent(
                        content_type=ContentType.HEADING,
                        text=text,
                        metadata={
                            "level": level,
                            "line": line_num + 1
                        },
                        position=current_position
                    ))
                
                # Extract lists
                elif line.startswith(('- ', '* ', '+ ')) or re.match(r'^\d+\.\s', line):
                    list_type = "ordered" if re.match(r'^\d+\.\s', line) else "unordered"
                    text = re.sub(r'^[-*+]\s*|\d+\.\s*', '', line)
                    
                    structured_content.append(ExtractedContent(
                        content_type=ContentType.LIST,
                        text=text,
                        metadata={
                            "list_type": list_type,
                            "line": line_num + 1
                        },
                        position=current_position
                    ))
                
                # Extract code blocks
                elif line.startswith('```'):
                    language = line[3:].strip()
                    code_lines = []
                    
                    for next_line_num in range(line_num + 1, len(lines)):
                        next_line = lines[next_line_num]
                        if next_line.strip() == '```':
                            break
                        code_lines.append(next_line)
                    
                    code_text = '\n'.join(code_lines)
                    
                    structured_content.append(ExtractedContent(
                        content_type=ContentType.CODE,
                        text=code_text,
                        metadata={
                            "language": language,
                            "start_line": line_num + 1,
                            "end_line": next_line_num + 1
                        },
                        position=current_position
                    ))
                
                # Extract links
                link_matches = re.finditer(r'\[([^\]]+)\]\(([^)]+)\)', line)
                for match in link_matches:
                    link_text = match.group(1)
                    link_url = match.group(2)
                    
                    document_structure.links.append({
                        "text": link_text,
                        "url": link_url,
                        "line": line_num + 1,
                        "type": "external" if link_url.startswith("http") else "internal"
                    })
                
                # Extract images
                image_matches = re.finditer(r'!\[([^\]]*)\]\(([^)]+)\)', line)
                for match in image_matches:
                    alt_text = match.group(1)
                    image_url = match.group(2)
                    
                    document_structure.images.append({
                        "alt_text": alt_text,
                        "url": image_url,
                        "line": line_num + 1
                    })
                
                # Regular paragraph
                else:
                    structured_content.append(ExtractedContent(
                        content_type=ContentType.PARAGRAPH,
                        text=line,
                        metadata={"line": line_num + 1},
                        position=current_position
                    ))
                
                current_position += len(line) + 1
            
            # Extract tables
            table_pattern = r'\|.*\|'
            table_lines = []
            in_table = False
            
            for line_num, line in enumerate(lines):
                if re.match(table_pattern, line.strip()):
                    if not in_table:
                        in_table = True
                        table_lines = []
                    table_lines.append(line.strip())
                else:
                    if in_table and table_lines:
                        # Process completed table
                        table_data = self._parse_markdown_table(table_lines)
                        document_structure.tables.append({
                            "start_line": line_num - len(table_lines) + 1,
                            "end_line": line_num,
                            "rows": len(table_data),
                            "columns": len(table_data[0]) if table_data else 0,
                            "data": table_data
                        })
                        
                        table_text = '\n'.join(table_lines)
                        structured_content.append(ExtractedContent(
                            content_type=ContentType.TABLE,
                            text=table_text,
                            metadata={
                                "start_line": line_num - len(table_lines) + 1,
                                "end_line": line_num,
                                "rows": len(table_data),
                                "columns": len(table_data[0]) if table_data else 0
                            }
                        ))
                    
                    in_table = False
            
            processing_time = time.time() - start_time
            
            result = ProcessingResult(
                document_id=document_id,
                format="markdown",
                status=ProcessingStatus.COMPLETED,
                extracted_text=extracted_text,
                structured_content=structured_content,
                document_structure=document_structure,
                metadata=metadata,
                processing_time=processing_time,
                confidence_score=1.0,
                word_count=len(extracted_text.split()),
                character_count=len(extracted_text)
            )
            
            DOCUMENT_PROCESSING_REQUESTS.labels(format='markdown', status='success').inc()
            CONTENT_EXTRACTION_ACCURACY.labels(format='markdown').set(1.0)
            
            return result
        
        except Exception as e:
            processing_time = time.time() - start_time
            logger.error(f"Markdown processing failed: {e}")
            
            DOCUMENT_PROCESSING_REQUESTS.labels(format='markdown', status='error').inc()
            
            return ProcessingResult(
                document_id=document_id,
                format="markdown",
                status=ProcessingStatus.FAILED,
                extracted_text="",
                structured_content=[],
                document_structure=DocumentStructure(),
                metadata={},
                processing_time=processing_time,
                error_message=str(e)
            )
        
        finally:
            latency = time.time() - start_time
            DOCUMENT_PROCESSING_LATENCY.labels(format='markdown').observe(latency)
    
    def _extract_front_matter(self, content: str) -> Dict[str, Any]:
        """Extract YAML front matter from markdown"""
        metadata = {}
        
        if content.startswith('---'):
            try:
                import yaml
                end_marker = content.find('\n---\n', 3)
                if end_marker != -1:
                    front_matter = content[3:end_marker]
                    metadata = yaml.safe_load(front_matter) or {}
            except ImportError:
                logger.warning("PyYAML not available for front matter parsing")
            except Exception as e:
                logger.warning(f"Failed to parse front matter: {e}")
        
        return metadata
    
    def _remove_front_matter(self, content: str) -> str:
        """Remove YAML front matter from markdown content"""
        if content.startswith('---'):
            end_marker = content.find('\n---\n', 3)
            if end_marker != -1:
                return content[end_marker + 5:]
        
        return content
    
    def _parse_markdown_table(self, table_lines: List[str]) -> List[List[str]]:
        """Parse markdown table into structured data"""
        table_data = []
        
        for line in table_lines:
            # Skip separator lines
            if re.match(r'^\|[\s\-:]+\|$', line.strip()):
                continue
            
            # Parse table row
            cells = [cell.strip() for cell in line.split('|')[1:-1]]
            if cells:
                table_data.append(cells)
        
        return table_data

class HTMLProcessor:
    """HTML document processor"""
    
    async def process_html(self, file_path: str) -> ProcessingResult:
        """Process HTML document"""
        start_time = time.time()
        document_id = hashlib.md5(file_path.encode()).hexdigest()
        
        try:
            if not BeautifulSoup:
                raise ImportError("beautifulsoup4 library not available")
            
            async with aiofiles.open(file_path, 'r', encoding='utf-8') as file:
                content = await file.read()
            
            soup = BeautifulSoup(content, 'html.parser')
            
            # Extract text content
            extracted_text = soup.get_text(separator='\n', strip=True)
            
            structured_content = []
            document_structure = DocumentStructure()
            
            # Extract metadata
            metadata = {}
            title_tag = soup.find('title')
            if title_tag:
                metadata['title'] = title_tag.get_text().strip()
            
            meta_tags = soup.find_all('meta')
            for meta in meta_tags:
                name = meta.get('name') or meta.get('property')
                content_attr = meta.get('content')
                if name and content_attr:
                    metadata[name] = content_attr
            
            # Extract headings
            for heading in soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6']):
                level = int(heading.name[1])
                text = heading.get_text().strip()
                
                document_structure.headings.append({
                    "level": level,
                    "text": text,
                    "tag": heading.name
                })
                
                structured_content.append(ExtractedContent(
                    content_type=ContentType.HEADING,
                    text=text,
                    metadata={
                        "level": level,
                        "tag": heading.name
                    }
                ))
            
            # Extract paragraphs
            for para in soup.find_all('p'):
                text = para.get_text().strip()
                if text:
                    structured_content.append(ExtractedContent(
                        content_type=ContentType.PARAGRAPH,
                        text=text,
                        metadata={"tag": "p"}
                    ))
            
            # Extract lists
            for list_elem in soup.find_all(['ul', 'ol']):
                list_type = "ordered" if list_elem.name == 'ol' else "unordered"
                items = []
                
                for li in list_elem.find_all('li'):
                    item_text = li.get_text().strip()
                    if item_text:
                        items.append(item_text)
                        
                        structured_content.append(ExtractedContent(
                            content_type=ContentType.LIST,
                            text=item_text,
                            metadata={
                                "list_type": list_type,
                                "tag": "li"
                            }
                        ))
            
            # Extract tables
            for table_index, table in enumerate(soup.find_all('table')):
                table_data = []
                
                for row in table.find_all('tr'):
                    row_data = []
                    for cell in row.find_all(['td', 'th']):
                        row_data.append(cell.get_text().strip())
                    if row_data:
                        table_data.append(row_data)
                
                document_structure.tables.append({
                    "index": table_index,
                    "rows": len(table_data),
                    "columns": len(table_data[0]) if table_data else 0,
                    "data": table_data
                })
                
                table_text = '\n'.join(['\t'.join(row) for row in table_data])
                structured_content.append(ExtractedContent(
                    content_type=ContentType.TABLE,
                    text=table_text,
                    metadata={
                        "table_index": table_index,
                        "rows": len(table_data),
                        "columns": len(table_data[0]) if table_data else 0
                    }
                ))
            
            # Extract links
            for link in soup.find_all('a', href=True):
                href = link['href']
                text = link.get_text().strip()
                
                document_structure.links.append({
                    "text": text,
                    "url": href,
                    "type": "external" if href.startswith("http") else "internal"
                })
            
            # Extract images
            for img in soup.find_all('img'):
                src = img.get('src', '')
                alt = img.get('alt', '')
                
                document_structure.images.append({
                    "src": src,
                    "alt": alt,
                    "width": img.get('width'),
                    "height": img.get('height')
                })
            
            # Extract code blocks
            for code in soup.find_all(['code', 'pre']):
                code_text = code.get_text()
                
                structured_content.append(ExtractedContent(
                    content_type=ContentType.CODE,
                    text=code_text,
                    metadata={"tag": code.name}
                ))
            
            processing_time = time.time() - start_time
            
            result = ProcessingResult(
                document_id=document_id,
                format="html",
                status=ProcessingStatus.COMPLETED,
                extracted_text=extracted_text,
                structured_content=structured_content,
                document_structure=document_structure,
                metadata=metadata,
                processing_time=processing_time,
                confidence_score=1.0,
                word_count=len(extracted_text.split()),
                character_count=len(extracted_text)
            )
            
            DOCUMENT_PROCESSING_REQUESTS.labels(format='html', status='success').inc()
            CONTENT_EXTRACTION_ACCURACY.labels(format='html').set(1.0)
            
            return result
        
        except Exception as e:
            processing_time = time.time() - start_time
            logger.error(f"HTML processing failed: {e}")
            
            DOCUMENT_PROCESSING_REQUESTS.labels(format='html', status='error').inc()
            
            return ProcessingResult(
                document_id=document_id,
                format="html",
                status=ProcessingStatus.FAILED,
                extracted_text="",
                structured_content=[],
                document_structure=DocumentStructure(),
                metadata={},
                processing_time=processing_time,
                error_message=str(e)
            )
        
        finally:
            latency = time.time() - start_time
            DOCUMENT_PROCESSING_LATENCY.labels(format='html').observe(latency)

class DocumentProcessor:
    """Main document processor that coordinates format-specific processors"""
    
    def __init__(self, config: Dict[str, Any]):
        self.config = config
        self.pdf_processor = PDFProcessor()
        self.office_processor = OfficeDocumentProcessor()
        self.markdown_processor = MarkdownProcessor()
        self.html_processor = HTMLProcessor()
        
        # Initialize NLP components if available
        self.nlp_enabled = self._initialize_nlp()
    
    def _initialize_nlp(self) -> bool:
        """Initialize NLP components"""
        try:
            if spacy:
                # Try to load English model
                try:
                    self.nlp = spacy.load("en_core_web_sm")
                    return True
                except OSError:
                    logger.warning("spaCy English model not found. NLP features disabled.")
            
            if nltk:
                # Download required NLTK data
                try:
                    nltk.download('punkt', quiet=True)
                    nltk.download('stopwords', quiet=True)
                    nltk.download('wordnet', quiet=True)
                    self.lemmatizer = WordNetLemmatizer()
                    return True
                except Exception as e:
                    logger.warning(f"NLTK initialization failed: {e}")
            
            return False
        
        except Exception as e:
            logger.warning(f"NLP initialization failed: {e}")
            return False
    
    async def process_document(self, file_path: str, format_hint: Optional[str] = None) -> ProcessingResult:
        """Process document based on format"""
        # Detect format if not provided
        if not format_hint:
            format_hint = self._detect_format(file_path)
        
        format_hint = format_hint.lower()
        
        # Route to appropriate processor
        if format_hint == "pdf":
            result = await self.pdf_processor.process_pdf(file_path)
        elif format_hint == "docx":
            result = await self.office_processor.process_docx(file_path)
        elif format_hint == "pptx":
            result = await self.office_processor.process_pptx(file_path)
        elif format_hint == "xlsx":
            result = await self.office_processor.process_xlsx(file_path)
        elif format_hint in ["markdown", "md"]:
            result = await self.markdown_processor.process_markdown(file_path)
        elif format_hint in ["html", "htm"]:
            result = await self.html_processor.process_html(file_path)
        elif format_hint == "txt":
            result = await self._process_text_file(file_path)
        else:
            result = ProcessingResult(
                document_id=hashlib.md5(file_path.encode()).hexdigest(),
                format=format_hint,
                status=ProcessingStatus.FAILED,
                extracted_text="",
                structured_content=[],
                document_structure=DocumentStructure(),
                metadata={},
                processing_time=0.0,
                error_message=f"Unsupported format: {format_hint}"
            )
        
        # Enhance with NLP if available
        if self.nlp_enabled and result.status == ProcessingStatus.COMPLETED:
            result = await self._enhance_with_nlp(result)
        
        return result
    
    def _detect_format(self, file_path: str) -> str:
        """Detect document format from file extension and MIME type"""
        # Get file extension
        _, ext = os.path.splitext(file_path)
        ext = ext.lower()
        
        # Map extensions to formats
        extension_map = {
            '.pdf': 'pdf',
            '.docx': 'docx',
            '.doc': 'docx',  # Treat as docx for now
            '.pptx': 'pptx',
            '.ppt': 'pptx',  # Treat as pptx for now
            '.xlsx': 'xlsx',
            '.xls': 'xlsx',  # Treat as xlsx for now
            '.md': 'markdown',
            '.markdown': 'markdown',
            '.html': 'html',
            '.htm': 'html',
            '.txt': 'txt',
            '.rtf': 'rtf'
        }
        
        if ext in extension_map:
            return extension_map[ext]
        
        # Try MIME type detection
        mime_type, _ = mimetypes.guess_type(file_path)
        if mime_type:
            if 'pdf' in mime_type:
                return 'pdf'
            elif 'word' in mime_type or 'officedocument.wordprocessingml' in mime_type:
                return 'docx'
            elif 'powerpoint' in mime_type or 'officedocument.presentationml' in mime_type:
                return 'pptx'
            elif 'excel' in mime_type or 'officedocument.spreadsheetml' in mime_type:
                return 'xlsx'
            elif 'html' in mime_type:
                return 'html'
            elif 'text' in mime_type:
                return 'txt'
        
        return 'unknown'
    
    async def _process_text_file(self, file_path: str) -> ProcessingResult:
        """Process plain text file"""
        start_time = time.time()
        document_id = hashlib.md5(file_path.encode()).hexdigest()
        
        try:
            async with aiofiles.open(file_path, 'r', encoding='utf-8') as file:
                content = await file.read()
            
            # Simple paragraph detection
            paragraphs = [p.strip() for p in content.split('\n\n') if p.strip()]
            
            structured_content = []
            for para in paragraphs:
                structured_content.append(ExtractedContent(
                    content_type=ContentType.PARAGRAPH,
                    text=para
                ))
            
            processing_time = time.time() - start_time
            
            result = ProcessingResult(
                document_id=document_id,
                format="txt",
                status=ProcessingStatus.COMPLETED,
                extracted_text=content,
                structured_content=structured_content,
                document_structure=DocumentStructure(),
                metadata={},
                processing_time=processing_time,
                confidence_score=1.0,
                word_count=len(content.split()),
                character_count=len(content)
            )
            
            DOCUMENT_PROCESSING_REQUESTS.labels(format='txt', status='success').inc()
            CONTENT_EXTRACTION_ACCURACY.labels(format='txt').set(1.0)
            
            return result
        
        except Exception as e:
            processing_time = time.time() - start_time
            logger.error(f"Text file processing failed: {e}")
            
            DOCUMENT_PROCESSING_REQUESTS.labels(format='txt', status='error').inc()
            
            return ProcessingResult(
                document_id=document_id,
                format="txt",
                status=ProcessingStatus.FAILED,
                extracted_text="",
                structured_content=[],
                document_structure=DocumentStructure(),
                metadata={},
                processing_time=processing_time,
                error_message=str(e)
            )
        
        finally:
            latency = time.time() - start_time
            DOCUMENT_PROCESSING_LATENCY.labels(format='txt').observe(latency)
    
    async def _enhance_with_nlp(self, result: ProcessingResult) -> ProcessingResult:
        """Enhance processing result with NLP analysis"""
        try:
            if spacy and hasattr(self, 'nlp'):
                # Process text with spaCy
                doc = self.nlp(result.extracted_text[:1000000])  # Limit to 1M chars
                
                # Extract entities
                entities = []
                for ent in doc.ents:
                    entities.append({
                        "text": ent.text,
                        "label": ent.label_,
                        "start": ent.start_char,
                        "end": ent.end_char,
                        "confidence": 0.9  # spaCy doesn't provide confidence scores
                    })
                
                # Detect language
                if hasattr(doc, 'lang_'):
                    result.language = doc.lang_
                
                # Add NLP metadata
                result.metadata['nlp'] = {
                    "entities": entities,
                    "sentences": len(list(doc.sents)),
                    "tokens": len(doc),
                    "language": result.language
                }
            
            elif nltk and hasattr(self, 'lemmatizer'):
                # Basic NLTK processing
                sentences = sent_tokenize(result.extracted_text)
                words = word_tokenize(result.extracted_text)
                
                # Remove stopwords and lemmatize
                try:
                    stop_words = set(stopwords.words('english'))
                    filtered_words = [self.lemmatizer.lemmatize(word.lower()) 
                                    for word in words 
                                    if word.lower() not in stop_words and word.isalpha()]
                    
                    result.metadata['nlp'] = {
                        "sentences": len(sentences),
                        "words": len(words),
                        "filtered_words": len(filtered_words),
                        "vocabulary": list(set(filtered_words))[:100]  # Top 100 unique words
                    }
                except Exception as e:
                    logger.warning(f"NLTK processing failed: {e}")
        
        except Exception as e:
            logger.warning(f"NLP enhancement failed: {e}")
        
        return result
    
    async def health_check(self) -> Dict[str, Any]:
        """Perform health check"""
        return {
            'status': 'healthy',
            'timestamp': datetime.utcnow().isoformat(),
            'processors': {
                'pdf': 'available' if PyPDF2 or fitz else 'unavailable',
                'office': 'available' if docx and pptx and openpyxl else 'partial',
                'markdown': 'available',
                'html': 'available' if BeautifulSoup else 'unavailable',
                'text': 'available'
            },
            'nlp': 'available' if self.nlp_enabled else 'unavailable',
            'ocr': 'available' if self.pdf_processor.ocr_enabled else 'unavailable'
        }

if __name__ == "__main__":
    import asyncio
    
    async def main():
        config = {}
        processor = DocumentProcessor(config)
        
        # Example usage
        test_files = [
            ("test.pdf", "pdf"),
            ("test.docx", "docx"),
            ("test.md", "markdown"),
            ("test.html", "html")
        ]
        
        for file_path, format_hint in test_files:
            if os.path.exists(file_path):
                result = await processor.process_document(file_path, format_hint)
                print(f"Processed {file_path}: {result.status.value}")
                print(f"  Extracted {result.word_count} words")
                print(f"  Processing time: {result.processing_time:.2f}s")
                print(f"  Confidence: {result.confidence_score:.2f}")
        
        # Health check
        health = await processor.health_check()
        print(f"Health: {health}")
    
    asyncio.run(main())

